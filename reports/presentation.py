# reports/presentation.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import pandas as pd
from datetime import datetime
import os

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "VoxSense"
    subtitle.text = "Voice Analysis & Classification System\nMachine Learning Project"

    # Slide 2: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Project Introduction"
    tf = content.text_frame
    tf.text = "• Real-time Audio Classification System\n"
    tf.add_paragraph().text = "• Uses MFCC, Chroma & Spectral Features\n"
    tf.add_paragraph().text = "• Deep Learning / Traditional ML Models\n"
    tf.add_paragraph().text = "• Built with Python, Librosa, PyTorch & Streamlit"

    # Slide 3: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Problem Statement"
    tf = content.text_frame
    tf.text = "Audio files ko automatically classify karna\n(Emotion / Speaker / Sound Type)"
    p = tf.add_paragraph()
    p.text = "Challenges:\n• Different audio formats\n• Noise & background sound\n• Limited dataset"

    # Slide 4: Project Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Project Architecture / Flow"
    tf = content.text_frame
    tf.text = "1. Data Collection (.mp3)\n"
    tf.add_paragraph().text = "2. Preprocessing (.mp3 → .wav)\n"
    tf.add_paragraph().text = "3. Feature Extraction (MFCC + Chroma)\n"
    tf.add_paragraph().text = "4. Spectrogram Generation\n"
    tf.add_paragraph().text = "5. Model Training (Neural Network)\n"
    tf.add_paragraph().text = "6. Evaluation & Deployment (Streamlit App)"

    # Slide 5: Technologies Used
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Technologies & Tools"
    tf = content.text_frame
    tf.text = "• Python 3\n"
    tf.add_paragraph().text = "• Librosa, SoundFile\n"
    tf.add_paragraph().text = "• PyTorch\n"
    tf.add_paragraph().text = "• Streamlit (Web Interface)\n"
    tf.add_paragraph().text = "• Scikit-learn\n"
    tf.add_paragraph().text = "• Matplotlib & Seaborn (Visualization)"

    # Slide 6: Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Results & Performance"
    tf = content.text_frame
    tf.text = "• Model Trained Successfully\n"
    tf.add_paragraph().text = "• Real-time Prediction via Web App\n"
    tf.add_paragraph().text = "• Performance Metrics & Graphs Generated\n"
    tf.add_paragraph().text = "• Confusion Matrix & Bar Charts"

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Conclusion & Future Work"
    tf = content.text_frame
    tf.text = "• Successfully built an end-to-end Voice Classification System\n"
    tf.add_paragraph().text = "• Future Scope:\n"
    tf.add_paragraph().text = "   - More classes (Emotions)\n"
    tf.add_paragraph().text = "   - CNN with Spectrograms\n"
    tf.add_paragraph().text = "   - Real-time microphone input"

    # Slide 8: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout for thank you
    title = slide.shapes.title
    title.text = "Thank You!"
    
    # Add date
    left = Inches(1)
    top = Inches(5)
    width = Inches(8)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = f"Presented on: {datetime.now().strftime('%d %B %Y')}"
    p.alignment = PP_ALIGN.CENTER

    # Save Presentation
    output_path = "reports/presentation.pptx"
    prs.save(output_path)

    print(f"🎉 Presentation created successfully!")
    print(f"📁 File saved at: {output_path}")
    print("\nAb aap ise directly open kar sakte ho.")

def main():
    create_presentation()

if __name__ == "__main__":
    main()